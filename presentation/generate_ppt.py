"""
Kanda Robot — Review I Presentation Generator
Run: python3 generate_ppt.py
Output: kanda_review1.pptx (same folder)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour Palette ────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x1B, 0x2A, 0x4A)
BLUE        = RGBColor(0x25, 0x63, 0xEB)
BLUE_LIGHT  = RGBColor(0xDB, 0xEA, 0xFE)
GREEN       = RGBColor(0x16, 0xA3, 0x4A)
GREEN_LIGHT = RGBColor(0xDC, 0xFC, 0xE7)
GRAY        = RGBColor(0x9C, 0xA3, 0xAF)
GRAY_LIGHT  = RGBColor(0xF3, 0xF4, 0xF6)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE      = RGBColor(0xF5, 0x9E, 0x0B)
RED_LIGHT   = RGBColor(0xFE, 0xF3, 0xC7)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height,
        fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, left, top, width, height,
        size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def header_bar(slide, title, subtitle=None):
    """Dark navy top bar with title."""
    box(slide, 0, 0, 13.33, 1.1, fill_color=NAVY)
    txt(slide, title,
        0.35, 0.12, 10, 0.6,
        size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle,
            0.35, 0.65, 10, 0.4,
            size=13, bold=False, color=BLUE_LIGHT, align=PP_ALIGN.LEFT, italic=True)
    # accent line
    acc = slide.shapes.add_shape(1, Inches(0), Inches(1.1), Inches(13.33), Inches(0.055))
    acc.fill.solid()
    acc.fill.fore_color.rgb = BLUE
    acc.line.fill.background()


def status_badge(slide, left, top, label, done=True):
    """Small green/gray pill badge."""
    color = GREEN if done else GRAY
    b = box(slide, left, top, 1.5, 0.28, fill_color=color)
    b.fill.solid()
    b.fill.fore_color.rgb = color
    b.line.fill.background()
    tf = b.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = WHITE


def bullet_block(slide, items, left, top, width, height,
                 heading=None, heading_color=BLUE,
                 bullet_color=NAVY, size=14, spacing=0.38):
    """Render a list of (text, done:bool|None) tuples.
       done=True → green dot, done=False → gray dot, done=None → blue dot"""
    y = top
    if heading:
        txt(slide, heading, left, y, width, 0.35,
            size=15, bold=True, color=heading_color)
        y += 0.38

    for (item, done) in items:
        dot_color = GREEN if done is True else (GRAY if done is False else BLUE)
        # bullet dot
        dot = slide.shapes.add_shape(9, Inches(left), Inches(y + 0.07),
                                     Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = dot_color
        dot.line.fill.background()

        item_color = GREEN if done is True else (GRAY if done is False else NAVY)
        txt(slide, item, left + 0.22, y, width - 0.22, spacing,
            size=size, color=item_color)
        y += spacing
    return y


def card(slide, left, top, width, height,
         title, body_lines, done=True, title_size=13, body_size=12):
    """A rounded-corner card with title + bullet lines."""
    bg_color = GREEN_LIGHT if done else GRAY_LIGHT
    border_color = GREEN if done else GRAY
    box(slide, left, top, width, height,
        fill_color=bg_color, line_color=border_color, line_width=Pt(1.2))
    # title
    title_bar = slide.shapes.add_shape(
        1, Inches(left), Inches(top),
        Inches(width), Inches(0.32)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = border_color
    title_bar.line.fill.background()
    tf = title_bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(title_size)
    run.font.bold = True
    run.font.color.rgb = WHITE

    y = top + 0.38
    for line in body_lines:
        txt(slide, f"• {line}", left + 0.1, y, width - 0.2, 0.32,
            size=body_size, color=NAVY if done else GRAY)
        y += 0.3


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1 — Title / Cover"""
    s = add_slide(prs)
    bg(s, NAVY)

    # gradient-ish blocks
    acc1 = s.shapes.add_shape(1, Inches(0), Inches(5.2), Inches(13.33), Inches(2.3))
    acc1.fill.solid(); acc1.fill.fore_color.rgb = RGBColor(0x11, 0x1B, 0x33)
    acc1.line.fill.background()

    blue_bar = s.shapes.add_shape(1, Inches(0), Inches(5.2), Inches(4.5), Inches(0.07))
    blue_bar.fill.solid(); blue_bar.fill.fore_color.rgb = BLUE
    blue_bar.line.fill.background()

    txt(s, "KANDA", 0.6, 0.7, 12, 1.4,
        size=72, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    txt(s, "A Multimodal Embodied Robot Agent", 0.65, 2.0, 12, 0.7,
        size=26, bold=False, color=BLUE_LIGHT, align=PP_ALIGN.LEFT)
    txt(s, "Powered by Large Language Models", 0.65, 2.6, 12, 0.7,
        size=26, bold=False, color=BLUE_LIGHT, align=PP_ALIGN.LEFT)

    txt(s, "Review I  ·  Phase 2: Embodiment Layer", 0.65, 5.5, 11, 0.45,
        size=15, bold=False, color=GRAY, align=PP_ALIGN.LEFT, italic=True)
    txt(s, "Hardware-Aware Autonomous Control System", 0.65, 5.9, 11, 0.4,
        size=13, color=GRAY, align=PP_ALIGN.LEFT)

    # phase pill
    status_badge(s, 10.8, 6.6, "PHASE 2  ✅", done=True)


def slide_agenda(prs):
    """Slide 2 — Agenda"""
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Agenda", "Review I  ·  Kanda Robot Project")

    items = [
        ("01  Introduction", None),
        ("02  Conceptual Understanding", None),
        ("03  Revisiting Objectives", None),
        ("04  Tool / Technique Selection & Justification", None),
        ("05  Solution Design Quality", None),
        ("06  Innovation & Problem-Solving", None),
        ("07  Partial Implementation", None),
        ("08  References", None),
    ]
    y = 1.4
    for i, (item, _) in enumerate(items):
        num, label = item.split("  ", 1)
        # number box
        nb = s.shapes.add_shape(1, Inches(0.55), Inches(y),
                                Inches(0.55), Inches(0.42))
        nb.fill.solid(); nb.fill.fore_color.rgb = BLUE
        nb.line.fill.background()
        nt = nb.text_frame; nt.word_wrap = False
        np_ = nt.paragraphs[0]; np_.alignment = PP_ALIGN.CENTER
        nr = np_.add_run(); nr.text = num
        nr.font.size = Pt(13); nr.font.bold = True; nr.font.color.rgb = WHITE

        txt(s, label, 1.25, y + 0.04, 11, 0.38, size=15, color=NAVY)
        y += 0.58


def slide_introduction(prs):
    """Slide 3 — Introduction"""
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Introduction", "What is Kanda?")

    txt(s,
        "Kanda is a hardware-grounded robotic agent designed to evolve "
        "from a rule-based autonomous controller into a fully LLM-driven "
        "embodied AI system — integrating physical sensing, motor control, "
        "and natural language reasoning.",
        0.5, 1.35, 12.3, 0.9, size=15, color=NAVY)

    # Two-column: Problem | Proposed Solution
    box(s, 0.5, 2.35, 5.9, 4.0, fill_color=BLUE_LIGHT,
        line_color=BLUE, line_width=Pt(1))
    txt(s, "Problem Gap", 0.55, 2.42, 5.7, 0.38,
        size=15, bold=True, color=BLUE)
    gap_items = [
        "Existing robots are purely rule-based",
        "No integration of AI reasoning with physical hardware",
        "Limited adaptability to novel environments",
        "Multimodal input (vision/audio) largely absent",
    ]
    y = 2.85
    for item in gap_items:
        txt(s, f"✗  {item}", 0.65, y, 5.6, 0.36, size=13, color=NAVY)
        y += 0.38

    box(s, 6.9, 2.35, 5.9, 4.0, fill_color=GREEN_LIGHT,
        line_color=GREEN, line_width=Pt(1))
    txt(s, "Proposed Solution", 6.95, 2.42, 5.7, 0.38,
        size=15, bold=True, color=GREEN)
    sol_items = [
        "Embodied robot with full sense–decide–act loop",
        "LLM as reasoning brain (Raspberry Pi bridge)",
        "Multimodal inputs: ultrasonic, camera, microphone",
        "Hardware-aware action generation via structured JSON",
        "Closed-loop feedback with real-time OLED display",
    ]
    y = 2.85
    for item in sol_items:
        txt(s, f"✓  {item}", 7.0, y, 5.7, 0.36, size=13, color=NAVY)
        y += 0.38


def slide_conceptual(prs):
    """Slide 4 — Conceptual Understanding"""
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Conceptual Understanding", "Embodied AI · Reactive vs Deliberative Agents")

    cols = [
        ("Embodied AI", BLUE, [
            "Intelligence that exists within a physical body",
            "Perception → Action loop is central",
            "Grounded in real-world sensor feedback",
            "Contrast: disembodied LLMs act on text only",
        ]),
        ("Reactive Agent\n(Current)", ORANGE, [
            "Responds directly to sensor inputs",
            "No internal world model",
            "Fast, deterministic, hardware-safe",
            "Limited — cannot generalize to new tasks",
        ]),
        ("Deliberative Agent\n(Target)", GREEN, [
            "Builds context before acting",
            "Uses LLM for reasoning & planning",
            "Can handle open-ended instructions",
            "Slower but far more capable",
        ]),
    ]

    x = 0.45
    for (title, color, items) in cols:
        fill = BLUE_LIGHT if color == BLUE else (RED_LIGHT if color == ORANGE else GREEN_LIGHT)
        box(s, x, 1.35, 3.9, 5.6, fill_color=fill,
            line_color=color, line_width=Pt(1.5))
        t_bar = s.shapes.add_shape(1, Inches(x), Inches(1.35), Inches(3.9), Inches(0.42))
        t_bar.fill.solid(); t_bar.fill.fore_color.rgb = color
        t_bar.line.fill.background()
        tf = t_bar.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = title
        r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE

        y = 1.88
        for item in items:
            txt(s, f"• {item}", x + 0.15, y, 3.65, 0.7,
                size=12.5, color=NAVY)
            y += 0.7
        x += 4.25


def slide_objectives(prs):
    """Slide 5 — Revisiting Objectives"""
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Revisiting Objectives", "Original Goals vs Current Status")

    rows = [
        ("Build a stable hardware platform",            "Design stable power + sensor + motor system",    True),
        ("Real-time obstacle detection",                 "3× HC-SR04 sensors with voltage dividers",       True),
        ("Autonomous navigation logic",                  "Sense→Decide→Act loop with smooth steering",     True),
        ("Real-time visual feedback",                    "SSD1306 OLED showing distances + decisions",     True),
        ("Clean, conflict-free pin mapping",             "Avoided strapping/UART/unstable GPIO pins",      True),
        ("Integrate LLM reasoning layer",                "Raspberry Pi bridge + LLM API pipeline",         False),
        ("Multimodal input (camera, mic)",               "Camera + microphone module integration",         False),
        ("Hardware-aware JSON action generation",        "Structured command protocol: Pi → ESP32",        False),
        ("Closed-loop LLM feedback",                     "Sensor telemetry fed back to LLM context",       False),
    ]

    # header row
    hx = [0.35, 3.45, 7.85, 11.2]
    for label, hx_ in zip(["#", "Objective", "Implementation Detail", "Status"], hx):
        box(s, hx_, 1.25, [0.7, 3.6, 3.6, 1.75][["#","Objective","Implementation Detail","Status"].index(label)],
            0.38, fill_color=NAVY)
        th = s.shapes[-1].text_frame; th.word_wrap = False
        hp = th.paragraphs[0]; hp.alignment = PP_ALIGN.CENTER
        hr = hp.add_run(); hr.text = label
        hr.font.size = Pt(11); hr.font.bold = True; hr.font.color.rgb = WHITE

    y = 1.65
    for i, (obj, detail, done) in enumerate(rows):
        row_bg = GREEN_LIGHT if done else GRAY_LIGHT
        box(s, 0.35, y, 12.6, 0.44, fill_color=row_bg,
            line_color=GREEN if done else GRAY, line_width=Pt(0.5))

        txt(s, str(i+1), 0.38, y + 0.08, 0.55, 0.32,
            size=11, bold=True, color=GREEN if done else GRAY, align=PP_ALIGN.CENTER)
        txt(s, obj, 3.5, y + 0.06, 3.5, 0.36, size=10.5, color=NAVY if done else GRAY)
        txt(s, detail, 7.9, y + 0.06, 3.5, 0.36, size=10.5, color=NAVY if done else GRAY)

        label = "✅ Done" if done else "⏳ Planned"
        lc = GREEN if done else GRAY
        txt(s, label, 11.25, y + 0.08, 1.6, 0.32,
            size=10.5, bold=True, color=lc, align=PP_ALIGN.CENTER)
        y += 0.46


def slide_tools(prs):
    """Slide 6 — Tool / Technique Selection"""
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Tool / Technique Selection & Justification", "Hardware + Software stack decisions")

    tools = [
        ("ESP32", "Microcontroller", "Dual-core, built-in WiFi/BT, 3.3V, abundant GPIOs, ledc PWM API", True),
        ("TB6612FNG", "Motor Driver", "Efficient H-bridge, 1.2A/ch, STBY pin, direction + PWM control", True),
        ("HC-SR04 ×3", "Ultrasonic Sensor", "Low cost, 2–400cm range; 1kΩ+1kΩ divider for 5V→2.5V ECHO", True),
        ("SSD1306 OLED", "Display", "I2C (2 pins), 128×64, real-time decision feedback", True),
        ("Buck Converter", "Power Regulation", "Stable 5V for ESP32; motors run on direct battery to avoid surge", True),
        ("Raspberry Pi", "AI Compute Bridge", "Full Linux OS, runs Python, calls LLM API, UART to ESP32", False),
        ("GPT-4 / Gemini", "LLM Reasoning", "Multimodal capability, function calling for structured JSON output", False),
        ("JSON Protocol", "Command Interface", "Type-safe, validatable, hardware-safe action specification", False),
    ]

    y = 1.38
    for (name, role, reason, done) in tools:
        row_bg = GREEN_LIGHT if done else GRAY_LIGHT
        border = GREEN if done else GRAY
        box(s, 0.35, y, 12.6, 0.52, fill_color=row_bg,
            line_color=border, line_width=Pt(0.6))

        # component name
        nb = s.shapes.add_shape(1, Inches(0.38), Inches(y + 0.07),
                                Inches(1.65), Inches(0.38))
        nb.fill.solid(); nb.fill.fore_color.rgb = border; nb.line.fill.background()
        nt = nb.text_frame; nt.word_wrap = False
        np_ = nt.paragraphs[0]; np_.alignment = PP_ALIGN.CENTER
        nr = np_.add_run(); nr.text = name
        nr.font.size = Pt(11); nr.font.bold = True; nr.font.color.rgb = WHITE

        txt(s, role, 2.18, y + 0.1, 2.0, 0.36, size=11,
            bold=True, color=NAVY if done else GRAY)
        txt(s, reason, 4.3, y + 0.1, 8.4, 0.36, size=11,
            color=NAVY if done else GRAY)
        y += 0.56


def slide_design_hld(prs):
    """Slide 7 — Solution Design: HLD"""
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Solution Design — High-Level Architecture (HLD)",
               "Full system view: completed (green) · planned (gray)")

    # Phase 2 box — LEFT
    box(s, 0.4, 1.35, 6.1, 5.65, fill_color=GREEN_LIGHT,
        line_color=GREEN, line_width=Pt(1.5))
    txt(s, "⚙️  Phase 2 — Embodiment Layer  ✅",
        0.5, 1.38, 5.9, 0.38, size=13, bold=True, color=GREEN)

    p2_components = [
        ("HC-SR04 (×3)", "Front · Left · Right ultrasonic"),
        ("ESP32 DevKit", "Central microcontroller"),
        ("TB6612FNG", "Dual motor driver"),
        ("SSD1306 OLED", "I2C real-time display"),
        ("LiPo + BMS + Buck", "Stable power architecture"),
        ("Left + Right Motor", "Physical locomotion"),
    ]
    y = 1.82
    for (comp, desc) in p2_components:
        box(s, 0.55, y, 5.7, 0.55, fill_color=GREEN,
            line_color=GREEN, line_width=Pt(0))
        b = s.shapes[-1].text_frame; b.word_wrap = False
        bp = b.paragraphs[0]; bp.alignment = PP_ALIGN.LEFT
        br = bp.add_run(); br.text = f"  {comp}"
        br.font.size = Pt(12); br.font.bold = True; br.font.color.rgb = WHITE

        txt(s, desc, 0.6, y + 0.28, 5.6, 0.25, size=10, color=WHITE,
            italic=True)

        # arrow between components
        if y < 4.7:
            arr = s.shapes.add_shape(1, Inches(2.8), Inches(y + 0.58),
                                     Inches(0.08), Inches(0.12))
            arr.fill.solid(); arr.fill.fore_color.rgb = GREEN; arr.line.fill.background()
        y += 0.72

    # Phase 3 box — RIGHT
    box(s, 6.9, 1.35, 6.0, 5.65, fill_color=GRAY_LIGHT,
        line_color=GRAY, line_width=Pt(1.5))
    txt(s, "🧠  Phase 3 — Intelligence Layer  ⏳",
        7.0, 1.38, 5.8, 0.38, size=13, bold=True, color=GRAY)

    p3_components = [
        ("Camera Module", "Visual perception input"),
        ("Microphone", "Audio / speech input"),
        ("Raspberry Pi", "AI compute bridge (brain)"),
        ("LLM API", "GPT-4 / Gemini reasoning"),
        ("Safety Validator", "Command verification layer"),
        ("Speaker (TTS)", "Audio output / feedback"),
    ]
    y = 1.82
    for (comp, desc) in p3_components:
        box(s, 7.05, y, 5.7, 0.55, fill_color=GRAY,
            line_color=GRAY, line_width=Pt(0))
        b = s.shapes[-1].text_frame; b.word_wrap = False
        bp = b.paragraphs[0]; bp.alignment = PP_ALIGN.LEFT
        br = bp.add_run(); br.text = f"  {comp}"
        br.font.size = Pt(12); br.font.bold = True; br.font.color.rgb = WHITE
        txt(s, desc, 7.1, y + 0.28, 5.6, 0.25, size=10, color=WHITE, italic=True)
        y += 0.72

    # Bridge arrow
    arr_box = s.shapes.add_shape(1, Inches(6.1), Inches(4.1), Inches(0.8), Inches(0.45))
    arr_box.fill.solid(); arr_box.fill.fore_color.rgb = BLUE; arr_box.line.fill.background()
    at = arr_box.text_frame; at.word_wrap = False
    ap = at.paragraphs[0]; ap.alignment = PP_ALIGN.CENTER
    ar = ap.add_run(); ar.text = "UART"
    ar.font.size = Pt(9); ar.font.bold = True; ar.font.color.rgb = WHITE


def slide_design_lld(prs):
    """Slide 8 — Solution Design: LLD (Pin Map)"""
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Solution Design — Low-Level Design (LLD)",
               "Pin mapping · Power architecture · Signal conditioning")

    # Pin table — left half
    txt(s, "ESP32 Pin Mapping", 0.4, 1.35, 5.5, 0.38,
        size=14, bold=True, color=BLUE)

    pin_rows = [
        ("GPIO 5",  "TRIG_F  (Front HC-SR04)",          True),
        ("GPIO 34", "ECHO_F  [input-only, 2.5V divider]", True),
        ("GPIO 13", "TRIG_L  (Left HC-SR04)",            True),
        ("GPIO 35", "ECHO_L  [input-only, 2.5V divider]", True),
        ("GPIO 4",  "TRIG_R  (Right HC-SR04)",           True),
        ("GPIO 32", "ECHO_R  [2.5V divider]",            True),
        ("GPIO 18", "AIN1 — Motor A dir",                True),
        ("GPIO 19", "AIN2 — Motor A dir",                True),
        ("GPIO 23", "PWMA   — Motor A speed",            True),
        ("GPIO 26", "BIN1 — Motor B dir",                True),
        ("GPIO 27", "BIN2 — Motor B dir",                True),
        ("GPIO 14", "PWMB   — Motor B speed",            True),
        ("GPIO 21", "SDA — OLED I2C",                    True),
        ("GPIO 22", "SCL — OLED I2C",                    True),
        ("GPIO 1/3","UART TX/RX → Raspberry Pi",         False),
    ]

    # header
    for lbl, lx, lw in [("GPIO", 0.4, 1.1), ("Function", 1.6, 4.3)]:
        box(s, lx, 1.76, lw, 0.3, fill_color=NAVY)
        ht = s.shapes[-1].text_frame; ht.word_wrap = False
        hp = ht.paragraphs[0]; hp.alignment = PP_ALIGN.CENTER
        hr = hp.add_run(); hr.text = lbl
        hr.font.size = Pt(10); hr.font.bold = True; hr.font.color.rgb = WHITE

    y = 2.08
    for (pin, func, done) in pin_rows:
        row_bg = GREEN_LIGHT if done else GRAY_LIGHT
        bc = GREEN if done else GRAY
        box(s, 0.4, y, 1.1, 0.3, fill_color=row_bg, line_color=bc, line_width=Pt(0.5))
        txt(s, pin, 0.42, y + 0.04, 1.05, 0.24, size=9.5, bold=True,
            color=GREEN if done else GRAY, align=PP_ALIGN.CENTER)
        box(s, 1.5, y, 4.4, 0.3, fill_color=row_bg, line_color=bc, line_width=Pt(0.5))
        txt(s, func, 1.55, y + 0.04, 4.3, 0.24, size=9.5,
            color=NAVY if done else GRAY)
        y += 0.305

    # Power arch — right half
    txt(s, "Power Architecture", 6.7, 1.35, 6.2, 0.38,
        size=14, bold=True, color=BLUE)

    pwr_nodes = [
        ("LiPo Battery  7.4V",  GREEN, 7.4, 1.82),
        ("BMS  (protection)",    GREEN, 7.4, 2.38),
        ("Main Switch",          GREEN, 7.4, 2.94),
        ("Buck Converter → 5V", GREEN, 7.4, 3.50),
        ("ESP32 VIN",            GREEN, 7.4, 4.06),
        ("TB6612FNG VM",         GREEN, 10.0, 3.20),
        ("Motors (L+R)",         GREEN, 10.0, 3.76),
        ("Raspberry Pi  ⏳",     GRAY,  10.0, 4.32),
    ]

    for (label, color, px, py) in pwr_nodes:
        b = s.shapes.add_shape(1, Inches(px), Inches(py), Inches(2.85), Inches(0.38))
        b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()
        bt = b.text_frame; bt.word_wrap = False
        bp = bt.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run(); br.text = label
        br.font.size = Pt(11); br.font.bold = True; br.font.color.rgb = WHITE

    # GND note
    txt(s, "★  All grounds tied together (common GND)",
        6.7, 4.62, 6.1, 0.38, size=11, color=NAVY, bold=True)
    txt(s, "★  Capacitors across motor rails — suppress voltage spikes",
        6.7, 4.98, 6.1, 0.38, size=11, color=NAVY)
    txt(s, "★  Motors NOT powered through buck (current surge protection)",
        6.7, 5.34, 6.1, 0.38, size=11, color=NAVY)


def slide_flow(prs):
    """Slide 9 — Flow Diagram (visual)"""
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Flow Diagram — Sense → Decide → Act",
               "Current: rule-based (green)  ·  Planned: LLM-driven (gray)")

    # ── CURRENT FLOW (left column) ──────────────────────────────────────────
    steps_current = [
        ("⚡ Power On & Init", GREEN),
        ("Read F / L / R Distances", GREEN),
        ("Front < 20 cm?", BLUE),
        ("Stop Motors · 200ms", GREEN),
        ("Left > Right?  →  Turn L / R", GREEN),
        ("Left < 15 cm?  →  Slight Right", GREEN),
        ("Right < 15 cm?  →  Slight Left", GREEN),
        ("FORWARD at speedVal", GREEN),
        ("Update OLED + Serial log", GREEN),
        ("Delay 100ms → loop", GREEN),
    ]

    txt(s, "Current Implementation  ✅", 0.4, 1.35, 5.8, 0.38,
        size=13, bold=True, color=GREEN)
    y = 1.82
    for (label, color) in steps_current:
        b = s.shapes.add_shape(1, Inches(0.45), Inches(y), Inches(5.7), Inches(0.42))
        b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()
        bt = b.text_frame; bt.word_wrap = False
        bp = bt.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run(); br.text = label
        br.font.size = Pt(11); br.font.bold = (color != BLUE); br.font.color.rgb = WHITE

        if y < 6.0:
            arr = s.shapes.add_shape(1, Inches(3.1), Inches(y + 0.42), Inches(0.1), Inches(0.14))
            arr.fill.solid(); arr.fill.fore_color.rgb = GREEN; arr.line.fill.background()
        y += 0.56

    # ── PLANNED FLOW (right column) ─────────────────────────────────────────
    steps_planned = [
        ("Camera + Mic + Sensor Input", GRAY),
        ("Build Multimodal Context", GRAY),
        ("LLM API Call (Raspberry Pi)", GRAY),
        ("Parse JSON Command", GRAY),
        ("Safety Validation", GRAY),
        ("Serial → ESP32", GRAY),
        ("Execute Motor Command", GRAY),
        ("Sensor telemetry → Pi (loop)", GRAY),
    ]

    txt(s, "Planned: LLM-Driven  ⏳", 7.0, 1.35, 5.8, 0.38,
        size=13, bold=True, color=GRAY)
    y = 1.82
    for (label, color) in steps_planned:
        b = s.shapes.add_shape(1, Inches(7.05), Inches(y), Inches(5.7), Inches(0.42))
        b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()
        bt = b.text_frame; bt.word_wrap = False
        bp = bt.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run(); br.text = label
        br.font.size = Pt(11); br.font.color.rgb = WHITE

        if y < 6.0:
            arr = s.shapes.add_shape(1, Inches(9.7), Inches(y + 0.42), Inches(0.1), Inches(0.14))
            arr.fill.solid(); arr.fill.fore_color.rgb = GRAY; arr.line.fill.background()
        y += 0.56

    # bridge label
    bridge = s.shapes.add_shape(1, Inches(6.1), Inches(3.9), Inches(0.85), Inches(0.55))
    bridge.fill.solid(); bridge.fill.fore_color.rgb = BLUE; bridge.line.fill.background()
    bt = bridge.text_frame; bt.word_wrap = True
    bp = bt.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run(); br.text = "Phase 3"
    br.font.size = Pt(10); br.font.bold = True; br.font.color.rgb = WHITE


def slide_usecase(prs):
    """Slide 10 — Use Case Diagram"""
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Use Case Diagram",
               "Actor → System interactions · completed (green) · planned (gray)")

    actors = [
        ("👤 User /\nEngineer",   1.0, 3.0, NAVY),
        ("🌍 Environment\n(Physical)", 1.0, 5.5, NAVY),
        ("🤖 LLM Service\n(API)",  11.5, 3.5, GRAY),
    ]

    for (name, ax, ay, color) in actors:
        # head
        h = s.shapes.add_shape(9, Inches(ax), Inches(ay),
                               Inches(0.38), Inches(0.38))
        h.fill.solid(); h.fill.fore_color.rgb = color; h.line.fill.background()
        txt(s, name, ax - 0.3, ay + 0.42, 1.2, 0.55,
            size=10, color=color, align=PP_ALIGN.CENTER)

    # System boundary
    sys_box = s.shapes.add_shape(1, Inches(2.5), Inches(1.35), Inches(8.8), Inches(5.65))
    sys_box.fill.solid(); sys_box.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFF)
    sys_box.line.color.rgb = BLUE; sys_box.line.width = Pt(1.5)
    txt(s, "Kanda Robot System", 6.2, 1.42, 4.0, 0.38,
        size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    use_cases = [
        # (label, x, y, done)
        ("Flash Firmware to ESP32",          3.0, 1.88, True),
        ("Monitor via Serial + OLED",        3.0, 2.48, True),
        ("Tune Speed / Threshold Params",    3.0, 3.08, True),
        ("Autonomous Obstacle Avoidance",    3.0, 3.68, True),
        ("Real-time Sensor Display",         3.0, 4.28, True),
        ("Give Voice / Text Command",        7.3, 1.88, False),
        ("Camera Scene Understanding",       7.3, 2.48, False),
        ("LLM Reasoning & Decision",         7.3, 3.08, False),
        ("Natural Language Action",          7.3, 3.68, False),
        ("Safety Override / Validation",     7.3, 4.28, False),
    ]

    for (label, ux, uy, done) in use_cases:
        color = GREEN if done else GRAY
        bg_c  = GREEN_LIGHT if done else GRAY_LIGHT
        uc_b = s.shapes.add_shape(1, Inches(ux), Inches(uy), Inches(3.7), Inches(0.42))
        uc_b.fill.solid(); uc_b.fill.fore_color.rgb = bg_c
        uc_b.line.color.rgb = color; uc_b.line.width = Pt(1)
        ut = uc_b.text_frame; ut.word_wrap = False
        up = ut.paragraphs[0]; up.alignment = PP_ALIGN.CENTER
        ur = up.add_run(); ur.text = label
        ur.font.size = Pt(10.5); ur.font.color.rgb = NAVY if done else GRAY

    # legend
    for i, (label, color) in enumerate([("✅ Completed", GREEN), ("⏳ Planned", GRAY)]):
        b = s.shapes.add_shape(1, Inches(0.4 + i*2.0), Inches(6.85),
                               Inches(1.7), Inches(0.3))
        b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()
        bt = b.text_frame; bt.word_wrap = False
        bp = bt.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run(); br.text = label
        br.font.size = Pt(10); br.font.bold = True; br.font.color.rgb = WHITE


def slide_innovation(prs):
    """Slide 11 — Innovation & Problem-Solving"""
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Innovation & Problem-Solving",
               "Engineering decisions that made the system work")

    innovations = [
        (
            "Hardware-First AI Foundation",
            GREEN,
            [
                "Built full embodiment layer before AI to ensure stable execution substrate",
                "Most academic projects bolt hardware onto AI — we invert this correctly",
                "Result: deterministic, safe motor control regardless of AI state",
            ],
        ),
        (
            "Voltage Divider Improvisation",
            ORANGE,
            [
                "HC-SR04 ECHO outputs 5V — unsafe for ESP32 3.3V GPIO",
                "No 2kΩ resistors available → used 1kΩ + 1kΩ (same ratio, same result)",
                "Demonstrates engineering principle: workable + safe > perfect",
            ],
        ),
        (
            "ESP32 Core v3 API Migration",
            BLUE,
            [
                "ledcSetup() deprecated in Arduino core v3 — caused compile errors",
                "Migrated to new ledcAttach(pin, freq, bits) single-call API",
                "Proactively future-proofed against further deprecation",
            ],
        ),
        (
            "Dual-Layer Architecture (Body/Brain)",
            GREEN if False else GRAY,
            [
                "ESP32 = Body: fast, deterministic hardware execution",
                "Raspberry Pi = Brain: LLM reasoning, context, planning",
                "Separation allows independent testing and safe fallback",
            ],
        ),
    ]

    x = 0.38
    row = 0
    for i, (title, color, points) in enumerate(innovations):
        col = i % 2
        if i == 2:
            x = 0.38
        cx = x + col * 6.5
        cy = 1.38 + row * 3.0

        fill = GREEN_LIGHT if color == GREEN else (
               RED_LIGHT if color == ORANGE else (
               BLUE_LIGHT if color == BLUE else GRAY_LIGHT))
        box(s, cx, cy, 6.2, 2.8, fill_color=fill,
            line_color=color, line_width=Pt(1.5))
        tb = s.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(6.2), Inches(0.4))
        tb.fill.solid(); tb.fill.fore_color.rgb = color; tb.line.fill.background()
        tf = tb.text_frame; tf.word_wrap = False
        tp = tf.paragraphs[0]; tp.alignment = PP_ALIGN.CENTER
        tr_ = tp.add_run(); tr_.text = title
        tr_.font.size = Pt(12); tr_.font.bold = True; tr_.font.color.rgb = WHITE

        py = cy + 0.5
        for pt in points:
            txt(s, f"• {pt}", cx + 0.15, py, 5.9, 0.62, size=11.5, color=NAVY)
            py += 0.64

        if col == 1:
            row += 1


def slide_implementation(prs):
    """Slide 12 — Partial Implementation"""
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Partial Implementation — Phase 2 Complete",
               "What is built · tested · and working now")

    # Left: done
    box(s, 0.38, 1.35, 6.1, 5.65, fill_color=GREEN_LIGHT,
        line_color=GREEN, line_width=Pt(1.5))
    txt(s, "✅  Completed & Working", 0.5, 1.42, 5.8, 0.38,
        size=13, bold=True, color=GREEN)

    done_items = [
        ("Power Architecture", "Battery → BMS → Buck → ESP32  (motors direct)"),
        ("3× HC-SR04 Sensors", "Voltage dividers on ECHO · GPIO34/35/32"),
        ("TB6612FNG Motor Driver", "AIN/BIN direction + PWMA/PWMB speed control"),
        ("SSD1306 OLED", "I2C on GPIO21/22 · real-time sensor + action display"),
        ("Obstacle Avoidance Logic", "Front threshold · side steering · smooth turning"),
        ("Pin Mapping", "Conflict-free · no strapping/UART/unstable pins"),
        ("ESP32 Core v3 Firmware", "ledcAttach() API · serial debug output"),
    ]
    y = 1.88
    for (comp, detail) in done_items:
        b = s.shapes.add_shape(1, Inches(0.52), Inches(y), Inches(5.8), Inches(0.62))
        b.fill.solid(); b.fill.fore_color.rgb = WHITE
        b.line.color.rgb = GREEN; b.line.width = Pt(0.7)
        txt(s, comp,   0.65, y + 0.04, 5.5, 0.28, size=11.5, bold=True, color=GREEN)
        txt(s, detail, 0.65, y + 0.32, 5.5, 0.28, size=10.5, color=NAVY)
        y += 0.72

    # Right: pending
    box(s, 6.85, 1.35, 6.1, 5.65, fill_color=GRAY_LIGHT,
        line_color=GRAY, line_width=Pt(1.5))
    txt(s, "⏳  Pending — Phase 3", 6.97, 1.42, 5.8, 0.38,
        size=13, bold=True, color=GRAY)

    pending_items = [
        ("Raspberry Pi Integration", "UART bridge ESP32↔Pi · JSON command protocol"),
        ("Camera Module", "Video frames → Pi → LLM context"),
        ("Microphone", "Audio input · speech-to-text pipeline"),
        ("LLM API Pipeline", "Prompt engineering · GPT-4/Gemini integration"),
        ("Safety Validator", "Command range checks before motor execution"),
        ("Speaker / TTS", "Audio feedback from robot"),
        ("Closed-Loop Telemetry", "Sensor data fed back into LLM context"),
    ]
    y = 1.88
    for (comp, detail) in pending_items:
        b = s.shapes.add_shape(1, Inches(6.98), Inches(y), Inches(5.8), Inches(0.62))
        b.fill.solid(); b.fill.fore_color.rgb = WHITE
        b.line.color.rgb = GRAY; b.line.width = Pt(0.7)
        txt(s, comp,   7.1, y + 0.04, 5.5, 0.28, size=11.5, bold=True, color=GRAY)
        txt(s, detail, 7.1, y + 0.32, 5.5, 0.28, size=10.5, color=GRAY)
        y += 0.72


def slide_references(prs):
    """Slide 13 — References"""
    s = add_slide(prs)
    bg(s, WHITE)
    header_bar(s, "References", "")

    refs = [
        ("[1]  Espressif Systems. (2024). ESP32 Technical Reference Manual. https://docs.espressif.com",
         "Hardware datasheet for ESP32 microcontroller"),
        ("[2]  Toshiba. (2020). TB6612FNG Dual DC Motor Driver Datasheet.",
         "Motor driver component specification"),
        ("[3]  JSR. HC-SR04 Ultrasonic Module User Guide.",
         "Ultrasonic sensor operational parameters"),
        ("[4]  Solomon Systech. SSD1306 OLED Datasheet.",
         "128×64 OLED display I2C communication protocol"),
        ("[5]  Brown, T., et al. (2020). Language Models are Few-Shot Learners. NeurIPS.",
         "Foundation paper on GPT-3 / LLM capabilities"),
        ("[6]  Ahn, M., et al. (2022). Do As I Can, Not As I Say: Grounding Language in Robotic Affordances. arXiv:2204.01691.",
         "SayCan — embodied AI with LLM reasoning"),
        ("[7]  Brohan, A., et al. (2023). RT-2: Vision-Language-Action Models Transfer Web Knowledge to Robotic Control. arXiv:2307.15818.",
         "Multimodal LLM-to-robot action transfer"),
        ("[8]  Peng, X., et al. (2023). EmbodiedGPT: Vision-Language Pre-Training via Embodied Chain of Thought. arXiv:2305.15021.",
         "Chain-of-thought reasoning for embodied agents"),
        ("[9]  Arduino. (2024). ESP32 Arduino Core v3 Migration Guide. https://docs.arduino.cc",
         "ledcAttach API changes from core v2 to v3"),
    ]

    y = 1.38
    for i, (ref, note) in enumerate(refs):
        even = i % 2 == 0
        box(s, 0.38, y, 12.55, 0.55,
            fill_color=BLUE_LIGHT if even else WHITE,
            line_color=BLUE, line_width=Pt(0.4))
        txt(s, ref,  0.5, y + 0.03, 12.2, 0.28, size=10, color=NAVY, bold=False)
        txt(s, note, 0.5, y + 0.30, 12.2, 0.22, size=9,  color=GRAY, italic=True)
        y += 0.59


def slide_thankyou(prs):
    """Slide 14 — Thank You"""
    s = add_slide(prs)
    bg(s, NAVY)

    acc = s.shapes.add_shape(1, Inches(0), Inches(5.8), Inches(13.33), Inches(1.7))
    acc.fill.solid(); acc.fill.fore_color.rgb = RGBColor(0x11, 0x1B, 0x33)
    acc.line.fill.background()

    bar = s.shapes.add_shape(1, Inches(0), Inches(5.8), Inches(5), Inches(0.07))
    bar.fill.solid(); bar.fill.fore_color.rgb = GREEN; bar.line.fill.background()

    txt(s, "Thank You", 0.6, 1.2, 12, 1.4,
        size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, "Kanda  ·  A Multimodal Embodied Robot Agent",
        0.6, 3.0, 12, 0.6, size=22, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
    txt(s, "Phase 2 Complete  ·  Foundation ready for LLM Integration",
        0.6, 3.65, 12, 0.45, size=15, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

    for i, (label, done) in enumerate([
        ("Hardware Platform ✅", True),
        ("Obstacle Avoidance ✅", True),
        ("OLED Feedback ✅", True),
        ("LLM Integration ⏳", False),
        ("Multimodal Input ⏳", False),
    ]):
        color = GREEN if done else GRAY
        b = s.shapes.add_shape(1, Inches(1.8 + i * 2.0), Inches(4.5),
                               Inches(1.8), Inches(0.45))
        b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()
        bt = b.text_frame; bt.word_wrap = False
        bp = bt.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run(); br.text = label
        br.font.size = Pt(10); br.font.bold = True; br.font.color.rgb = WHITE


# ═════════════════════════════════════════════════════════════════════════════
# MAIN
# ═════════════════════════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_agenda(prs)
    slide_introduction(prs)
    slide_conceptual(prs)
    slide_objectives(prs)
    slide_tools(prs)
    slide_design_hld(prs)
    slide_design_lld(prs)
    slide_flow(prs)
    slide_usecase(prs)
    slide_innovation(prs)
    slide_implementation(prs)
    slide_references(prs)
    slide_thankyou(prs)

    out = "kanda_review1.pptx"
    prs.save(out)
    print(f"✅  Saved → {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
