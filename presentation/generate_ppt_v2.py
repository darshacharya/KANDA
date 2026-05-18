"""
KANDA Presentation Generator v2 — with real images
Run: python3 generate_ppt_v2.py
Output: kanda_review1_v2.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Image paths ───────────────────────────────────────────────────────────────
IMGS = os.path.join(os.path.dirname(__file__), "..", "imgs")
IMG_ROBOT1  = os.path.join(IMGS, "kanda(1).jpeg")   # top-down view
IMG_ROBOT2  = os.path.join(IMGS, "kanda(2).jpeg")   # front/sensor view
IMG_ARCH    = os.path.join(IMGS, "system-architecture.png")
IMG_LLD     = os.path.join(IMGS, "pin-level-lld.png")
IMG_FLOW    = os.path.join(IMGS, "flowchart.png")

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x1B, 0x2A, 0x4A)
BLUE        = RGBColor(0x25, 0x63, 0xEB)
BLUE_LIGHT  = RGBColor(0xDB, 0xEA, 0xFE)
GREEN       = RGBColor(0x16, 0xA3, 0x4A)
GREEN_LIGHT = RGBColor(0xDC, 0xFC, 0xE7)
GRAY        = RGBColor(0x9C, 0xA3, 0xAF)
GRAY_LIGHT  = RGBColor(0xF3, 0xF4, 0xF6)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE      = RGBColor(0xF5, 0x9E, 0x0B)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, l, t, w, h, fill_color=None, line_color=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill_color:
        s.fill.solid(); s.fill.fore_color.rgb = fill_color
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color; s.line.width = lw
    else:
        s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h, size=14, bold=False, color=NAVY,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb


def header(slide, title, subtitle=None):
    box(slide, 0, 0, 13.33, 1.1, fill_color=NAVY)
    acc = slide.shapes.add_shape(1, Inches(0), Inches(1.1), Inches(13.33), Inches(0.055))
    acc.fill.solid(); acc.fill.fore_color.rgb = BLUE; acc.line.fill.background()
    txt(slide, title, 0.35, 0.12, 10, 0.6, size=28, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.35, 0.67, 12, 0.38, size=12, color=BLUE_LIGHT, italic=True)


def img(slide, path, l, t, w, h):
    """Add image, maintain aspect ratio within given box."""
    return slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))


def caption(slide, text, l, t, w):
    b = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(0.3))
    b.fill.solid(); b.fill.fore_color.rgb = NAVY; b.line.fill.background()
    tf = b.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = WHITE


def bullet(slide, items, l, t, w, size=13, spacing=0.42):
    y = t
    for (text, done) in items:
        dot = slide.shapes.add_shape(9, Inches(l), Inches(y + 0.1),
                                     Inches(0.11), Inches(0.11))
        dot.fill.solid()
        dot.fill.fore_color.rgb = GREEN if done is True else (GRAY if done is False else BLUE)
        dot.line.fill.background()
        c = GREEN if done is True else (GRAY if done is False else NAVY)
        txt(slide, text, l + 0.2, y, w - 0.2, spacing, size=size, color=c)
        y += spacing
    return y


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    s = add_slide(prs)
    bg(s, NAVY)

    # Robot photo as right-side background
    img(s, IMG_ROBOT2, 7.5, 0.5, 5.5, 6.5)

    # Dark overlay on photo side for readability
    ov = s.shapes.add_shape(1, Inches(7.3), Inches(0), Inches(6.03), Inches(7.5))
    ov.fill.solid(); ov.fill.fore_color.rgb = RGBColor(0x0A, 0x10, 0x22)
    from pptx.util import Pt as uPt
    from pptx.dml.color import RGBColor as RGB
    ov.fill.fore_color.rgb = RGBColor(0x1B, 0x2A, 0x4A)
    ov.line.fill.background()
    # set transparency via XML
    sp_pr = ov._element.spPr
    from lxml import etree
    solid_fill = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    if solid_fill is not None:
        srgb = solid_fill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgb is None:
            srgb = solid_fill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}sysClr')
        if srgb is not None:
            alpha = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha.set('val', '75000')  # 75% opacity

    # Left content
    acc = s.shapes.add_shape(1, Inches(0), Inches(5.3), Inches(7.3), Inches(0.07))
    acc.fill.solid(); acc.fill.fore_color.rgb = GREEN; acc.line.fill.background()

    txt(s, "KANDA", 0.55, 0.7, 7, 1.5, size=72, bold=True, color=WHITE)
    txt(s, "Knowledge-driven Autonomous Navigation", 0.6, 2.1, 7, 0.55,
        size=16, color=BLUE_LIGHT)
    txt(s, "and Decision-making Agent", 0.6, 2.62, 7, 0.45, size=16, color=BLUE_LIGHT)
    txt(s, "A Multimodal Embodied Robot Agent Powered by LLMs",
        0.6, 3.22, 7, 0.38, size=13, bold=True, color=WHITE)
    txt(s, "with Hardware Aware Action Generation",
        0.6, 3.60, 7, 0.38, size=13, bold=True, color=WHITE)

    txt(s, "Review I  ·  Phase 2: Embodiment Layer Complete",
        0.6, 5.55, 7, 0.38, size=12, color=GRAY, italic=True)

    tag = s.shapes.add_shape(1, Inches(0.55), Inches(6.1), Inches(2.0), Inches(0.38))
    tag.fill.solid(); tag.fill.fore_color.rgb = GREEN; tag.line.fill.background()
    tf = tag.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "PHASE 2  ✅  COMPLETE"
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = WHITE


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — TABLE OF CONTENTS
# ═════════════════════════════════════════════════════════════════════════════

def slide_toc(prs):
    s = add_slide(prs)
    bg(s, WHITE)
    header(s, "Table of Contents", "Review I  ·  KANDA Robot Project")

    items = [
        "01   Introduction",
        "02   Meet KANDA",
        "03   System Design",
        "04   Methodology",
        "05   Use of Tools & Techniques",
        "06   Implementation Progress",
        "07   Identification of Journal / Conference",
        "08   References",
    ]
    y = 1.4
    for item in items:
        num, label = item.split("   ", 1)
        nb = s.shapes.add_shape(1, Inches(0.55), Inches(y), Inches(0.55), Inches(0.42))
        nb.fill.solid(); nb.fill.fore_color.rgb = BLUE; nb.line.fill.background()
        nt = nb.text_frame; nt.word_wrap = False
        np_ = nt.paragraphs[0]; np_.alignment = PP_ALIGN.CENTER
        nr = np_.add_run(); nr.text = num
        nr.font.size = Pt(13); nr.font.bold = True; nr.font.color.rgb = WHITE
        txt(s, label, 1.25, y + 0.04, 11, 0.38, size=15, color=NAVY)
        y += 0.58


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — INTRODUCTION
# ═════════════════════════════════════════════════════════════════════════════

def slide_intro(prs):
    s = add_slide(prs)
    bg(s, WHITE)
    header(s, "Introduction", "A Multimodal Embodied Robot Agent Powered by LLMs with Hardware Aware Action Generation")

    txt(s,
        "KANDA is a low-cost, embodied robot designed for household companion applications — "
        "combining physical sensing, autonomous navigation, and LLM-based reasoning on "
        "commodity hardware (ESP32 + Raspberry Pi).",
        0.5, 1.32, 12.3, 0.55, size=13.5, color=NAVY)

    # Left: 4 points
    box(s, 0.5, 2.0, 5.85, 5.0, fill_color=GREEN_LIGHT, line_color=GREEN, lw=Pt(1.2))
    hb = s.shapes.add_shape(1, Inches(0.5), Inches(2.0), Inches(5.85), Inches(0.42))
    hb.fill.solid(); hb.fill.fore_color.rgb = GREEN; hb.line.fill.background()
    ht = hb.text_frame; ht.word_wrap = False
    hp = ht.paragraphs[0]; hp.alignment = PP_ALIGN.CENTER
    hr = hp.add_run(); hr.text = "What KANDA Does"
    hr.font.size = Pt(13); hr.font.bold = True; hr.font.color.rgb = WHITE

    points = [
        ("Autonomous navigation on commodity hardware — ESP32 + HC-SR04 sensors + motor driver",
         True),
        ("LLM as the reasoning brain — Raspberry Pi bridges robot to GPT-4 / Gemini API",
         False),
        ("Household companion use cases — elderly care, child tutoring, reminders, home automation",
         None),
        ("Safety-first LLM-to-hardware execution — every command validated before motor execution",
         False),
    ]
    y = 2.52
    for (text, done) in points:
        dot_color = GREEN if done is True else (GRAY if done is False else BLUE)
        dot = s.shapes.add_shape(9, Inches(0.68), Inches(y + 0.08),
                                 Inches(0.13), Inches(0.13))
        dot.fill.solid(); dot.fill.fore_color.rgb = dot_color; dot.line.fill.background()
        txt(s, text, 0.9, y, 5.3, 0.82, size=12, color=NAVY if done is True else
            (GRAY if done is False else NAVY))
        y += 0.88

    # Right: research gap
    box(s, 6.98, 2.0, 5.85, 5.0, fill_color=BLUE_LIGHT, line_color=BLUE, lw=Pt(1.2))
    hb2 = s.shapes.add_shape(1, Inches(6.98), Inches(2.0), Inches(5.85), Inches(0.42))
    hb2.fill.solid(); hb2.fill.fore_color.rgb = BLUE; hb2.line.fill.background()
    ht2 = hb2.text_frame; ht2.word_wrap = False
    hp2 = ht2.paragraphs[0]; hp2.alignment = PP_ALIGN.CENTER
    hr2 = hp2.add_run(); hr2.text = "Research Gap Addressed"
    hr2.font.size = Pt(13); hr2.font.bold = True; hr2.font.color.rgb = WHITE

    gap_points = [
        "LLM-robot systems (SayCan, RT-2) require expensive proprietary hardware",
        "Low-cost ESP32 systems (ESP-SparkBot) have voice only — no navigation",
        "No open protocol translates LLM output into hardware-safe motor commands",
        "Companion robotics research uses commercial robots, not DIY platforms",
    ]
    y = 2.55
    for gp in gap_points:
        txt(s, f"✗  {gp}", 7.1, y, 5.55, 0.95, size=12, color=NAVY)
        y += 0.95


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — MEET KANDA (robot photos)
# ═════════════════════════════════════════════════════════════════════════════

def slide_meet_kanda(prs):
    s = add_slide(prs)
    bg(s, WHITE)
    header(s, "Meet KANDA", "Phase 2 — Built & Working  ✅")

    # Photo 1 — top view (left)
    img(s, IMG_ROBOT1, 0.4, 1.25, 5.8, 5.1)
    caption(s, "Top View — ESP32 · TB6612FNG · OLED · Battery", 0.4, 6.38, 5.8)

    # Photo 2 — front/sensor view (right)
    img(s, IMG_ROBOT2, 6.8, 1.25, 6.1, 5.1)
    caption(s, "Front View — HC-SR04 Sensors · Wheels · Full Assembly", 6.8, 6.38, 6.1)

    # Component labels (left photo callouts)
    for (label, color, lx, ly) in [
        ("Battery + BMS", GREEN, 0.5, 1.32),
        ("ESP32", GREEN, 0.5, 1.7),
        ("TB6612FNG", GREEN, 0.5, 2.08),
        ("OLED Display", GREEN, 0.5, 2.46),
        ("Breadboard wiring", GREEN, 0.5, 2.84),
    ]:
        b = s.shapes.add_shape(1, Inches(lx), Inches(ly), Inches(1.55), Inches(0.3))
        b.fill.solid(); b.fill.fore_color.rgb = GREEN
        b.line.fill.background()
        bt = b.text_frame; bt.word_wrap = False
        bp = bt.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run(); br.text = label
        br.font.size = Pt(8.5); br.font.bold = True; br.font.color.rgb = WHITE


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — SYSTEM DESIGN (architecture diagram)
# ═════════════════════════════════════════════════════════════════════════════

def slide_system_design(prs):
    s = add_slide(prs)
    bg(s, WHITE)
    header(s, "System Design — Architecture",
           "Two-layer design: Embodiment (✅ done) · Intelligence (⏳ planned)")

    # Architecture diagram — main area
    img(s, IMG_ARCH, 0.4, 1.25, 12.5, 5.5)

    # Legend
    for i, (label, color) in enumerate([("✅ Completed", GREEN), ("⏳ Planned", GRAY)]):
        b = s.shapes.add_shape(1, Inches(0.4 + i * 2.1), Inches(6.85),
                               Inches(1.85), Inches(0.3))
        b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()
        bt = b.text_frame; bt.word_wrap = False
        bp = bt.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run(); br.text = label
        br.font.size = Pt(10); br.font.bold = True; br.font.color.rgb = WHITE


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — LLD (pin-level diagram)
# ═════════════════════════════════════════════════════════════════════════════

def slide_lld(prs):
    s = add_slide(prs)
    bg(s, WHITE)
    header(s, "Low-Level Design — Pin Mapping & Wiring",
           "Signal conditioning · Power architecture · Component connections")

    img(s, IMG_LLD, 0.4, 1.25, 12.5, 5.6)

    # key notes strip at bottom
    notes = [
        "ECHO pins: 5V → 2.5V via 1kΩ+1kΩ divider",
        "PWM: 1kHz · 8-bit · ledcAttach() API",
        "I2C: GPIO21 SDA · GPIO22 SCL",
        "Motors: direct battery feed (not through buck)",
    ]
    x = 0.4
    for note in notes:
        b = s.shapes.add_shape(1, Inches(x), Inches(6.88), Inches(3.0), Inches(0.28))
        b.fill.solid(); b.fill.fore_color.rgb = NAVY; b.line.fill.background()
        bt = b.text_frame; bt.word_wrap = False
        bp = bt.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run(); br.text = note
        br.font.size = Pt(9); br.font.color.rgb = WHITE
        x += 3.15


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — METHODOLOGY (flow diagram)
# ═════════════════════════════════════════════════════════════════════════════

def slide_methodology(prs):
    s = add_slide(prs)
    bg(s, WHITE)
    header(s, "Methodology — Sense → Decide → Act Loop",
           "Current: rule-based (green)  ·  Planned: LLM-driven (gray)")

    # Flow diagram — left 2/3
    img(s, IMG_FLOW, 0.4, 1.25, 8.2, 5.85)

    # Phase summary — right 1/3
    box(s, 9.0, 1.25, 4.0, 2.7, fill_color=GREEN_LIGHT, line_color=GREEN, lw=Pt(1.2))
    txt(s, "Phase 2  ✅  Complete", 9.1, 1.3, 3.8, 0.38, size=12, bold=True, color=GREEN)
    p2 = [
        "Read F/L/R sensor distances",
        "Rule-based obstacle logic",
        "Motor control via PWM",
        "OLED + Serial feedback",
    ]
    y = 1.75
    for p in p2:
        txt(s, f"• {p}", 9.15, y, 3.75, 0.36, size=11.5, color=NAVY)
        y += 0.38

    box(s, 9.0, 4.1, 4.0, 2.95, fill_color=GRAY_LIGHT, line_color=GRAY, lw=Pt(1.2))
    txt(s, "Phase 3  ⏳  Planned", 9.1, 4.15, 3.8, 0.38, size=12, bold=True, color=GRAY)
    p3 = [
        "Camera + Mic + Sensor input",
        "LLM API call (Raspberry Pi)",
        "Safety validator",
        "JSON → UART → ESP32",
        "TTS verbal response",
    ]
    y = 4.58
    for p in p3:
        txt(s, f"• {p}", 9.15, y, 3.75, 0.36, size=11.5, color=GRAY)
        y += 0.38


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — TOOLS & TECHNIQUES
# ═════════════════════════════════════════════════════════════════════════════

def slide_tools(prs):
    s = add_slide(prs)
    bg(s, WHITE)
    header(s, "Use of Tools & Techniques",
           "Hardware + Software stack with justification")

    tools = [
        ("ESP32",          "Microcontroller", "Dual-core, 3.3V, abundant GPIO, ledcAttach PWM API",           True),
        ("TB6612FNG",      "Motor Driver",    "Efficient H-bridge, 1.2A/ch, STBY pin, direction+PWM control", True),
        ("HC-SR04 ×3",     "Ultrasonic",      "2–400cm range; 1kΩ+1kΩ divider → 5V ECHO to 2.5V safe",       True),
        ("SSD1306 OLED",   "Display",         "I2C (2 pins), 128×64, real-time decision feedback",             True),
        ("Buck Converter",  "Power Reg.",      "Stable 5V for ESP32; motors run on direct battery to avoid surge", True),
        ("Raspberry Pi",   "AI Compute",      "Full Linux OS, Python, LLM API client, UART to ESP32",          False),
        ("GPT-4 / Gemini", "LLM Reasoning",   "Multimodal capability, JSON function calling for motor commands", False),
        ("Whisper (OpenAI)","Speech-to-Text",  "Converts microphone audio to text for LLM context",             False),
        ("JSON Protocol",  "Command Bridge",  "Type-safe, range-validated, hardware-safe action interface",     False),
    ]

    y = 1.35
    for (name, role, reason, done) in tools:
        bg_c = GREEN_LIGHT if done else GRAY_LIGHT
        border = GREEN if done else GRAY
        box(s, 0.35, y, 12.6, 0.52, fill_color=bg_c, line_color=border, lw=Pt(0.6))
        nb = s.shapes.add_shape(1, Inches(0.38), Inches(y + 0.07),
                                Inches(1.65), Inches(0.38))
        nb.fill.solid(); nb.fill.fore_color.rgb = border; nb.line.fill.background()
        nt = nb.text_frame; nt.word_wrap = False
        np_ = nt.paragraphs[0]; np_.alignment = PP_ALIGN.CENTER
        nr = np_.add_run(); nr.text = name
        nr.font.size = Pt(11); nr.font.bold = True; nr.font.color.rgb = WHITE

        txt(s, role,   2.18, y + 0.1, 2.0, 0.36, size=11, bold=True,
            color=NAVY if done else GRAY)
        txt(s, reason, 4.35, y + 0.1, 8.4, 0.36, size=11,
            color=NAVY if done else GRAY)
        y += 0.56


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — IMPLEMENTATION PROGRESS
# ═════════════════════════════════════════════════════════════════════════════

def slide_implementation(prs):
    s = add_slide(prs)
    bg(s, WHITE)
    header(s, "Implementation Progress",
           "Phase 2 complete · Phase 3 defined and ready to build")

    # Left: done items
    box(s, 0.35, 1.3, 6.1, 4.4, fill_color=GREEN_LIGHT, line_color=GREEN, lw=Pt(1.5))
    txt(s, "✅  Phase 2 — Complete", 0.48, 1.36, 5.8, 0.38, size=13, bold=True, color=GREEN)

    done_items = [
        "Power architecture: Battery → BMS → Buck → ESP32",
        "3× HC-SR04 sensors with voltage dividers (5V→2.5V)",
        "TB6612FNG motor control: direction + PWM speed",
        "SSD1306 OLED: real-time sensor + decision display",
        "Conflict-free pin mapping (no strapping/UART pins)",
        "Obstacle avoidance firmware (ESP32 Arduino core v3)",
        "Smooth steering: front stop + side correction logic",
    ]
    y = 1.82
    for item in done_items:
        b = s.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(5.8), Inches(0.46))
        b.fill.solid(); b.fill.fore_color.rgb = WHITE
        b.line.color.rgb = GREEN; b.line.width = Pt(0.6)
        txt(s, f"✓  {item}", 0.65, y + 0.08, 5.5, 0.3, size=11, color=GREEN)
        y += 0.5

    # Right: pending items
    box(s, 6.85, 1.3, 6.1, 4.4, fill_color=GRAY_LIGHT, line_color=GRAY, lw=Pt(1.5))
    txt(s, "⏳  Phase 3 — Planned", 6.98, 1.36, 5.8, 0.38, size=13, bold=True, color=GRAY)

    pending_items = [
        "Raspberry Pi integration + UART bridge to ESP32",
        "Camera module: visual scene input to LLM context",
        "Microphone + Whisper speech transcription pipeline",
        "LLM API integration (GPT-4 / Gemini)",
        "Safety validator: range checks before motor execution",
        "Text-to-speech verbal companion responses",
        "Closed-loop: sensor telemetry back to LLM context",
    ]
    y = 1.82
    for item in pending_items:
        b = s.shapes.add_shape(1, Inches(6.98), Inches(y), Inches(5.8), Inches(0.46))
        b.fill.solid(); b.fill.fore_color.rgb = WHITE
        b.line.color.rgb = GRAY; b.line.width = Pt(0.6)
        txt(s, f"○  {item}", 7.13, y + 0.08, 5.5, 0.3, size=11, color=GRAY)
        y += 0.5

    # Robot photo strip at bottom
    img(s, IMG_ROBOT2, 0.35, 5.82, 4.0, 1.4)
    img(s, IMG_ROBOT1, 4.55, 5.82, 4.0, 1.4)
    caption(s, "KANDA — Phase 2 Hardware (Built & Tested)", 0.35, 7.1, 8.2)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — JOURNAL / CONFERENCE
# ═════════════════════════════════════════════════════════════════════════════

def slide_journal(prs):
    s = add_slide(prs)
    bg(s, WHITE)
    header(s, "Identification of Journal / Conference",
           "Target venues for publishing KANDA research")

    venues = [
        ("IEEE Robotics and Automation Letters (RA-L)",
         "Q1 · Impact Factor ~5.2",
         "Top venue for embodied AI + robot control research. Directly relevant to hardware-aware LLM action generation.",
         True),
        ("ACM/IEEE HRI — Human-Robot Interaction Conference",
         "Top-tier conference",
         "Premier venue for companion robotics research. Covers elderly assistance, child interaction, and social robots.",
         True),
        ("Frontiers in Robotics and AI",
         "Q2 · Open Access",
         "Published companion robot design studies (2024). Strong fit for KANDA's household application domain.",
         True),
        ("JMIR Human Factors",
         "Q1 · Open Access",
         "Published LLM+robot geriatric care studies (2025). Direct match for elderly companion evaluation results.",
         False),
        ("IEEE Access",
         "Q2 · Open Access",
         "Broad engineering scope. Suitable for the embedded systems + AI integration aspect of KANDA.",
         False),
    ]

    y = 1.35
    for (name, tag, desc, primary) in venues:
        bg_c = GREEN_LIGHT if primary else GRAY_LIGHT
        border = GREEN if primary else GRAY
        box(s, 0.35, y, 12.6, 0.9, fill_color=bg_c, line_color=border, lw=Pt(0.8))

        nb = s.shapes.add_shape(1, Inches(0.38), Inches(y + 0.08),
                                Inches(4.5), Inches(0.35))
        nb.fill.solid(); nb.fill.fore_color.rgb = border; nb.line.fill.background()
        nt = nb.text_frame; nt.word_wrap = False
        np_ = nt.paragraphs[0]; np_.alignment = PP_ALIGN.LEFT
        nr = np_.add_run(); nr.text = f"  {name}"
        nr.font.size = Pt(11.5); nr.font.bold = True; nr.font.color.rgb = WHITE

        txt(s, tag,  5.05, y + 0.1,  2.5, 0.3, size=11, bold=True,
            color=GREEN if primary else GRAY)
        txt(s, desc, 7.7,  y + 0.1,  5.1, 0.7, size=11,
            color=NAVY if primary else GRAY)

        lbl = "Primary Target" if primary else "Secondary Option"
        lb = s.shapes.add_shape(1, Inches(5.05), Inches(y + 0.48),
                                Inches(2.4), Inches(0.26))
        lb.fill.solid(); lb.fill.fore_color.rgb = border; lb.line.fill.background()
        lt = lb.text_frame; lt.word_wrap = False
        lp = lt.paragraphs[0]; lp.alignment = PP_ALIGN.CENTER
        lr = lp.add_run(); lr.text = lbl
        lr.font.size = Pt(9); lr.font.bold = True; lr.font.color.rgb = WHITE

        y += 1.0


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — REFERENCES
# ═════════════════════════════════════════════════════════════════════════════

def slide_references(prs):
    s = add_slide(prs)
    bg(s, WHITE)
    header(s, "References", "")

    refs = [
        "[1]  M. Ahn et al., \"Do As I Can, Not As I Say,\" arXiv:2204.01691, Google, 2022.",
        "[2]  A. Brohan et al., \"RT-2: Vision-Language-Action Models,\" arXiv:2307.15818, Google DeepMind, 2023.",
        "[3]  D. Driess et al., \"PaLM-E: An Embodied Multimodal Language Model,\" arXiv:2303.03378, 2023.",
        "[4]  \"Safety Guardrails for LLM-Enabled Robots (RoboGuard),\" arXiv:2503.07885, 2025.",
        "[5]  \"SafeEmbodAI: Safety Validation for Mobile Robot Navigation,\" arXiv:2409.01630, 2024.",
        "[6]  J. Shen et al., \"Socially Assistive Robot + LLM in Geriatric Care,\" JMIR Human Factors, 2025.",
        "[7]  \"Conversational Companion Robots With Older Adults,\" Frontiers in Robotics and AI, 2024.",
        "[8]  Espressif Systems, \"ESP-SparkBot: LLM Robot with ESP32-S3,\" Espressif Dev Portal, 2025.",
        "[9]  A. Radford et al., \"Robust Speech Recognition via Large-Scale Weak Supervision (Whisper),\" OpenAI, 2022.",
        "[10] Espressif Systems, ESP32 Technical Reference Manual, v5.2, 2024.",
        "[11] Toshiba Semiconductor, TB6612FNG Motor Driver Datasheet, 2020.",
        "[12] Solomon Systech, SSD1306 OLED Driver Datasheet, 2008.",
    ]

    y = 1.38
    for i, ref in enumerate(refs):
        even = i % 2 == 0
        box(s, 0.38, y, 12.55, 0.42,
            fill_color=BLUE_LIGHT if even else WHITE,
            line_color=BLUE, lw=Pt(0.4))
        txt(s, ref, 0.5, y + 0.06, 12.2, 0.32, size=10.5, color=NAVY)
        y += 0.46


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — THANK YOU
# ═════════════════════════════════════════════════════════════════════════════

def slide_thankyou(prs):
    s = add_slide(prs)
    bg(s, NAVY)

    img(s, IMG_ROBOT2, 8.5, 1.0, 4.5, 5.5)

    acc = s.shapes.add_shape(1, Inches(0), Inches(5.85), Inches(13.33), Inches(1.65))
    acc.fill.solid(); acc.fill.fore_color.rgb = RGBColor(0x11, 0x1B, 0x33)
    acc.line.fill.background()

    bar = s.shapes.add_shape(1, Inches(0), Inches(5.85), Inches(4.5), Inches(0.07))
    bar.fill.solid(); bar.fill.fore_color.rgb = GREEN; bar.line.fill.background()

    txt(s, "Thank You", 0.5, 1.2, 8, 1.4, size=60, bold=True, color=WHITE)
    txt(s, "A Multimodal Embodied Robot Agent Powered by LLMs with Hardware Aware Action Generation",
        0.5, 2.9, 8, 0.7, size=14, color=BLUE_LIGHT)
    txt(s, "Phase 2 Complete · Foundation ready for LLM Integration",
        0.5, 3.55, 7, 0.42, size=13, color=GRAY, italic=True)

    for i, (label, done) in enumerate([
        ("Hardware ✅", True), ("Navigation ✅", True),
        ("OLED ✅", True), ("LLM ⏳", False), ("Companion ⏳", False),
    ]):
        b = s.shapes.add_shape(1, Inches(0.5 + i * 1.6), Inches(4.6),
                               Inches(1.45), Inches(0.42))
        b.fill.solid(); b.fill.fore_color.rgb = GREEN if done else GRAY
        b.line.fill.background()
        bt = b.text_frame; bt.word_wrap = False
        bp = bt.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run(); br.text = label
        br.font.size = Pt(10); br.font.bold = True; br.font.color.rgb = WHITE

    txt(s, '"The hardware works. The architecture is sound. KANDA is ready to think."',
        0.5, 6.1, 8.0, 0.5, size=12, color=WHITE, italic=True)


# ═════════════════════════════════════════════════════════════════════════════
# BUILD
# ═════════════════════════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_toc(prs)
    slide_intro(prs)
    slide_meet_kanda(prs)
    slide_system_design(prs)
    slide_lld(prs)
    slide_methodology(prs)
    slide_tools(prs)
    slide_implementation(prs)
    slide_journal(prs)
    slide_references(prs)
    slide_thankyou(prs)

    out = "kanda_review1_v2.pptx"
    prs.save(out)
    print(f"✅  Saved → {out}  ({len(prs.slides)} slides)")

if __name__ == "__main__":
    build()
